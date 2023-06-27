from pptx import Presentation


def extract_text_from_presentation(presentation_path):
    # Load the PowerPoint presentation
    presentation = Presentation(presentation_path)
    all_text = []  # List to store the extracted text from each slide

    # Iterate over each slide in the presentation
    for slide in presentation.slides:
        slide_text = []  # List to store the extracted text from the current slide

        # Iterate over each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Iterate over each paragraph in the shape's text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Iterate over each run in the paragraph
                    for run in paragraph.runs:
                        slide_text.append(run.text)  # Append the text to slide_text list

        # Join the extracted text from the current slide into a single string
        slide_text_combined = '\n'.join(slide_text)
        all_text.append(slide_text_combined)  # Append the text to all_text list

    return all_text
